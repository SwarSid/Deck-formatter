"""
Deck Formatting Agent — Streamlit App
======================================
Deploy to Streamlit Cloud for a shareable link your whole team can use.

Steps:
  1. Create free GitHub account at github.com
  2. New repo 'deck-formatter' — upload app.py + requirements.txt
  3. Go to share.streamlit.io → New app → pick repo → Deploy
  4. Share the URL with your team. Done.

requirements.txt:
  streamlit
  python-pptx
  lxml
  scipy
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree
import re, io, math

# ─────────────────────────────────────────────────────────────────
#  PAGE CONFIG
# ─────────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="Deck Formatter",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
  .block-container { padding-top:2rem; padding-bottom:2rem; }
  div[data-testid="stSidebar"] { background:#F8FAFF; }
  .section-label {
    font-size:11px; font-weight:700; letter-spacing:1.5px;
    text-transform:uppercase; color:#2563EB;
    margin-bottom:5px; margin-top:4px;
  }
  .fix-item    { padding:6px 10px; border-radius:6px; background:#F0F7FF;
                 border-left:3px solid #2563EB; font-size:13px; margin-bottom:4px; }
  .statsig-item{ padding:6px 10px; border-radius:6px; background:#F0FDF4;
                 border-left:3px solid #16A34A; font-size:13px; margin-bottom:4px; }
  .overlap-item{ padding:6px 10px; border-radius:6px; background:#FFF7ED;
                 border-left:3px solid #D97706; font-size:13px; margin-bottom:4px; }
  .warn-item   { padding:6px 10px; border-radius:6px; background:#FEF2F2;
                 border-left:3px solid #DC2626; font-size:13px; margin-bottom:4px; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────
#  XML NAMESPACES
# ─────────────────────────────────────────────────────────────────

_NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ─────────────────────────────────────────────────────────────────
#  SIDEBAR — ALL CHECKLIST SECTIONS
# ─────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("## 🎯 Deck Formatter")
    st.caption("Fill in your client's rules. Upload your deck on the right. Click Run.")
    st.markdown("---")

    # ── TYPOGRAPHY ────────────────────────────────────────────────
    st.markdown('<div class="section-label">🔤 Typography</div>', unsafe_allow_html=True)
    global_font = st.text_input("Font for all elements", value="Century Gothic",
        help="Applied to all text — chart labels, footnotes, axis labels")

    st.markdown('<div class="section-label">Talking Header</div>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    header_size      = c1.number_input("Size (pt)", value=20, min_value=6, max_value=72, key="hs")
    header_color_hex = c2.color_picker("Colour", value="#001F60", key="hc")

    st.markdown('<div class="section-label">Sub Header</div>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    subheader_size       = c1.number_input("Sub header (pt)", value=11, min_value=6, max_value=48, key="ss")
    below_subheader_size = c2.number_input("Below sub header (pt)", value=7, min_value=6, max_value=48, key="bs")

    st.markdown('<div class="section-label">Chart Title</div>', unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    chart_title_size      = c1.number_input("Size (pt)", value=11, min_value=6, max_value=48, key="cts")
    chart_title_bold      = c2.selectbox("Bold?", ["Yes", "No"], key="ctb")
    chart_title_color_hex = c3.color_picker("Colour", value="#000000", key="ctc")

    st.markdown("---")

    # ── CHART DATA LABELS ─────────────────────────────────────────
    st.markdown('<div class="section-label">📊 Chart Data Labels</div>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    label_min_pct        = c1.number_input("Remove labels below (%)", value=5,   min_value=0, max_value=50,  key="lmp")
    label_dark_threshold = c2.number_input("Dark bar threshold (0–255)", value=128, min_value=0, max_value=255, key="ldt")
    c1, c2 = st.columns(2)
    label_dark_color_hex  = c1.color_picker("Label on dark bars",  value="#FFFFFF", key="ldc")
    label_light_color_hex = c2.color_picker("Label on light bars", value="#000000", key="llc")
    label_remove_highlight = st.toggle("Remove highlights from data labels", value=True, key="lrh")

    st.markdown("---")

    # ── OVERLAP PREVENTION ────────────────────────────────────────
    st.markdown('<div class="section-label">↕ Label Overlap Prevention</div>', unsafe_allow_html=True)
    st.caption("Nudge labels apart when two products have close values (e.g. 76% vs 73%).")
    overlap_enabled = st.toggle("Enable overlap fix", value=True, key="oe")
    if overlap_enabled:
        c1, c2 = st.columns(2)
        overlap_threshold = c1.number_input("Threshold (%)", value=8,  min_value=1, max_value=30, key="ot")
        nudge_pts         = c2.number_input("Nudge (pts)",   value=14, min_value=4, max_value=40, key="np")
    else:
        overlap_threshold, nudge_pts = 8, 14

    st.markdown("---")

    # ── FOOTNOTES ─────────────────────────────────────────────────
    st.markdown('<div class="section-label">📝 Footnotes</div>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    footnote_size        = c1.number_input("Size (pt)", value=8, min_value=6, max_value=24, key="fs")
    footnote_color_hex   = c2.color_picker("Footnote colour", value="#808080", key="fc")
    c1, c2 = st.columns(2)
    footnote_lowsample_hex  = c1.color_picker("Low sample colour", value="#FF0000", key="flsc")
    footnote_keywords_str   = c2.text_input("Low sample keywords", value="low sample, s =, n =, *")
    footnote_remove_highlight = st.toggle("Remove highlights from footnotes", value=True, key="frh")

    st.markdown("---")

    # ── PLACEMENT ─────────────────────────────────────────────────
    st.markdown('<div class="section-label">📐 Header Placement</div>', unsafe_allow_html=True)
    st.caption("Snap all headers to the same position as Slide 1.")
    placement_enabled = st.toggle("Enforce consistent placement", value=True, key="pe")
    if placement_enabled:
        align_talking   = st.checkbox("Align talking headers", value=True, key="at")
        align_subheader = st.checkbox("Align sub headers",     value=True, key="ash")
        align_chart     = st.checkbox("Align chart titles",    value=True, key="ac")
        placement_mode  = st.selectbox("Alignment mode",
            ["Top position only (recommended)", "Top + Left position", "Left position only"], key="pm")
    else:
        align_talking = align_subheader = align_chart = False
        placement_mode = "Top position only (recommended)"

    st.markdown("---")

    # ── STAT SIG TESTING ──────────────────────────────────────────
    st.markdown('<div class="section-label">🟢 Statistical Significance</div>', unsafe_allow_html=True)
    st.caption("Adds a green ● next to labels that are statistically significant at 95% confidence.")

    statsig_enabled = st.toggle("Enable stat sig testing", value=False, key="sse")

    if statsig_enabled:

        st.markdown("**Which slides to test?**")
        statsig_slides_str = st.text_input(
            "Slide numbers (e.g. 2, 4, 7) — leave blank for all",
            value="", key="sss",
        )

        st.markdown("**Comparison mode**")
        st.caption(
            "• **Products vs Others** — tests each named product against the 'Others' series\n"
            "• **Current vs Previous quarter** — tests the most recent wave against the one before it"
        )
        statsig_mode = st.selectbox(
            "What to compare",
            ["Products vs Others", "Current vs Previous quarter"],
            key="sm",
        )

        st.markdown("**How sample sizes are found**")
        st.caption(
            "The script automatically reads **s = N** or **n = N** values from "
            "category axis labels or footnotes on each slide. "
            "No manual entry needed. If it cannot find an n, it will flag that in the log."
        )

        st.markdown("**Confidence level**")
        confidence_level = st.selectbox("Confidence level", ["95%", "99%", "90%"], key="cl")
        z_critical = {"95%": 1.96, "99%": 2.576, "90%": 1.645}[confidence_level]

        st.markdown("**Marker symbol**")
        st.caption("Choose the symbol that appears next to stat sig labels.")

        SYMBOL_OPTIONS = {
            "● Filled circle":    "●",
            "▲ Triangle":         "▲",
            "★ Star":             "★",
            "■ Square":           "■",
            "◆ Diamond":          "◆",
            "✦ Four-point star":  "✦",
            "* Asterisk":         "*",
            "↑ Up arrow":         "↑",
        }
        symbol_choice = st.selectbox(
            "Symbol",
            list(SYMBOL_OPTIONS.keys()),
            key="sym",
        )
        statsig_marker = SYMBOL_OPTIONS[symbol_choice]

        st.markdown("**Marker colour**")
        st.caption("Pick from standard options or use a custom colour.")

        # Named colour palette — covers the most common brand palettes
        COLOUR_OPTIONS = {
            "🟢 Green (standard stat sig)":   "#16A34A",
            "🔵 Navy blue":                   "#001F5B",
            "🔷 Royal blue":                   "#2563EB",
            "🟦 Teal":                         "#0D9488",
            "🟠 Amber / orange":               "#D97706",
            "🔴 Red":                          "#DC2626",
            "🟣 Purple":                       "#7C3AED",
            "⚫ Black":                         "#000000",
            "⚪ White (for dark backgrounds)": "#FFFFFF",
            "🎨 Custom colour…":               "custom",
        }

        colour_choice = st.selectbox(
            "Colour",
            list(COLOUR_OPTIONS.keys()),
            key="colour_choice",
        )

        if COLOUR_OPTIONS[colour_choice] == "custom":
            statsig_marker_color_hex = st.color_picker(
                "Pick your custom colour", value="#16A34A", key="smc_custom"
            )
        else:
            statsig_marker_color_hex = COLOUR_OPTIONS[colour_choice]
            # Show a preview of symbol + colour
            preview_hex = statsig_marker_color_hex
            st.markdown(
                f"<div style='font-size:22px;color:{preview_hex};"
                f"background:#F8FAFF;padding:6px 12px;"
                f"border-radius:8px;border:1px solid #E2E8F0;"
                f"display:inline-block;margin-top:4px'>"
                f"{statsig_marker} Preview</div>",
                unsafe_allow_html=True,
            )

        st.caption(
            "**Products vs Others**: tests each product series against the series whose "
            "name contains 'other' or 'oth'.\n\n"
            "**Current vs Previous quarter**: detects quarter labels (Q1/Q2/FY etc.) and "
            "compares the two most recent adjacent time points."
        )

    else:
        statsig_slides_str       = ""
        statsig_mode             = "Products vs Others"
        z_critical               = 1.96
        statsig_marker           = "●"
        statsig_marker_color_hex = "#16A34A"


# ─────────────────────────────────────────────────────────────────
#  HELPERS — shared
# ─────────────────────────────────────────────────────────────────

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def parse_slide_list(s):
    s = s.strip()
    if not s:
        return None
    nums = set()
    for p in s.split(","):
        p = p.strip()
        if p.isdigit():
            nums.add(int(p))
    return nums or None

def rgb_color(triple):
    return RGBColor(*triple)

def luminance(r, g, b):
    return 0.299*r + 0.587*g + 0.114*b

def is_dark_bar(r, g, b):
    return luminance(r, g, b) < CHECKLIST["data_label_dark_threshold"]

def remove_highlight(run):
    rPr = run._r.get_or_add_rPr()
    hl  = rPr.find(f"{{{_NS_A}}}highlight")
    if hl is not None:
        rPr.remove(hl)
    try:
        run.font.fill.background()
    except Exception:
        pass

def set_run(run, size=None, bold=None, color=None, font=None):
    if font:            run.font.name      = font
    if size:            run.font.size      = Pt(size)
    if bold is not None: run.font.bold     = bold
    if color:           run.font.color.rgb = rgb_color(color)

def pct_value(text):
    m = re.search(r"(\d+(?:\.\d+)?)\s*%", str(text))
    return float(m.group(1)) if m else None

def is_low_sample(text):
    t = text.lower()
    return any(k.lower() in t for k in CHECKLIST["footnote_low_sample_keywords"])

def is_footnote(shape, h):
    return shape.has_text_frame and (shape.top or 0) > h * 0.82

def is_header(shape, h):
    return shape.has_text_frame and (shape.top or 0) < h * 0.20

def is_subheader(shape, h):
    top = shape.top or 0
    return shape.has_text_frame and h * 0.20 <= top < h * 0.40

# ─────────────────────────────────────────────────────────────────
#  BUILD CHECKLIST
# ─────────────────────────────────────────────────────────────────

STATSIG_SLIDES = parse_slide_list(statsig_slides_str) if statsig_enabled else None

CHECKLIST = {
    "global_font":               global_font,
    "talking_header_size":       header_size,
    "talking_header_color":      hex_to_rgb(header_color_hex),
    "subheader_size":            subheader_size,
    "below_subheader_size":      below_subheader_size,
    "chart_title_size":          chart_title_size,
    "chart_title_bold":          chart_title_bold == "Yes",
    "chart_title_color":         hex_to_rgb(chart_title_color_hex),
    "data_label_min_pct":        label_min_pct,
    "data_label_dark_threshold": label_dark_threshold,
    "data_label_light_color":    hex_to_rgb(label_light_color_hex),
    "data_label_dark_color":     hex_to_rgb(label_dark_color_hex),
    "data_label_remove_highlight": label_remove_highlight,
    "label_overlap_enabled":     overlap_enabled,
    "label_overlap_threshold":   overlap_threshold,
    "label_nudge_pts":           nudge_pts,
    "footnote_size":             footnote_size,
    "footnote_color":            hex_to_rgb(footnote_color_hex),
    "footnote_low_sample_color": hex_to_rgb(footnote_lowsample_hex),
    "footnote_remove_highlight": footnote_remove_highlight,
    "footnote_low_sample_keywords": [k.strip() for k in footnote_keywords_str.split(",")],
    "statsig_enabled":           statsig_enabled,
    "statsig_slides":            STATSIG_SLIDES,
    "statsig_mode":              statsig_mode if statsig_enabled else "Products vs Others",
    "statsig_z_critical":        z_critical,
    "statsig_marker":            statsig_marker if statsig_enabled else "●",
    "statsig_marker_color":      hex_to_rgb(statsig_marker_color_hex) if statsig_enabled else (22, 163, 74),
}

# ─────────────────────────────────────────────────────────────────
#  STAT SIG ENGINE
# ─────────────────────────────────────────────────────────────────

def _extract_n_from_text(text):
    """
    Extract sample size from strings like:
      's = 75', 'n = 30', '(s=75)', 'S=29*', 'n=54'
    Returns int or None.
    """
    m = re.search(r"[sSnN]\s*=\s*(\d+)", str(text))
    return int(m.group(1)) if m else None

def _extract_all_n_from_slide(slide):
    """
    Scan all text shapes on the slide for sample size annotations.
    Returns dict: {label_text_fragment: n}
    e.g. {'FY26 Q2': 75, 'FY26 Q1': 75, 'NeuroOncs+': 21}
    """
    found = {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = para.text.strip()
            n = _extract_n_from_text(full)
            if n:
                # Store under the full text and any recognisable label fragment
                found[full] = n
                # Try to extract a label key (text before the s= part)
                label_part = re.split(r"\s*[\(\[]?\s*[sSnN]\s*=", full)[0].strip()
                if label_part:
                    found[label_part] = n
    return found

def _two_prop_z_test(p1, n1, p2, n2):
    """
    Two-proportion z-test.
    p1, p2 are proportions (0–100 scale accepted, converted internally).
    Returns z-score. Caller compares abs(z) against z_critical.
    """
    # Convert percentages to proportions if needed
    if p1 > 1: p1 /= 100.0
    if p2 > 1: p2 /= 100.0

    # Pooled proportion
    p_pool = (p1 * n1 + p2 * n2) / (n1 + n2)
    denom  = math.sqrt(p_pool * (1 - p_pool) * (1/n1 + 1/n2))
    if denom == 0:
        return 0.0
    return abs(p1 - p2) / denom

def _is_quarter_label(text):
    """Detect time-series labels: FY26 Q2, Q1 2024, Wave 3, etc."""
    t = text.upper()
    return bool(re.search(r"\bQ[1-4]\b|\bFY\d\d\b|\bWAVE\s*\d\b|\b20\d\d\b", t))

def _is_others_series(name):
    """Detect 'Others', 'Oth', 'Other' series names."""
    return bool(re.search(r"\bother", str(name).lower()))

def _get_series_values(chart):
    """
    Returns list of dicts:
      [{ 'name': str, 'values': [float|None, ...], 'series_obj': ... }, ...]
    One entry per series in chart order.
    """
    ns = _NS_C
    result = []
    try:
        plot_el = chart._element.find(f".//{{{ns}}}plotArea")
        if plot_el is None:
            return result
        SUPPORTED = {"barChart","bar3DChart","lineChart","line3DChart","areaChart","area3DChart"}
        for ctype in plot_el:
            if etree.QName(ctype.tag).localname not in SUPPORTED:
                continue
            for ser_el in ctype.findall(f"{{{ns}}}ser"):
                # Series name
                ser_name = ""
                tx_el = ser_el.find(f".//{{{ns}}}ser/{{{ns}}}tx/{{{ns}}}strRef/{{{ns}}}strCache/{{{ns}}}pt/{{{ns}}}v")
                if tx_el is None:
                    tx_el = ser_el.find(f".//{{{ns}}}tx/{{{ns}}}strRef/{{{ns}}}strCache/{{{ns}}}pt/{{{ns}}}v")
                if tx_el is None:
                    tx_el = ser_el.find(f".//{{{ns}}}tx/{{{ns}}}v")
                if tx_el is not None:
                    ser_name = tx_el.text or ""

                # Values
                cache = ser_el.find(f".//{{{ns}}}val/{{{ns}}}numRef/{{{ns}}}numCache")
                if cache is None:
                    cache = ser_el.find(f".//{{{ns}}}numRef/{{{ns}}}numCache")
                values = []
                if cache is not None:
                    pt_count = 0
                    for pt in cache.findall(f"{{{ns}}}pt"):
                        idx = int(pt.get("idx", pt_count))
                        v_el = pt.find(f"{{{ns}}}v")
                        val  = float(v_el.text) if v_el is not None and v_el.text else None
                        # Pad to idx
                        while len(values) <= idx:
                            values.append(None)
                        values[idx] = val
                        pt_count += 1

                # Category labels
                cat_labels = []
                cat_cache = ser_el.find(f".//{{{ns}}}cat/{{{ns}}}strRef/{{{ns}}}strCache")
                if cat_cache is None:
                    cat_cache = chart._element.find(
                        f".//{{{ns}}}cat/{{{ns}}}strRef/{{{ns}}}strCache")
                if cat_cache is not None:
                    for pt in cat_cache.findall(f"{{{ns}}}pt"):
                        v = pt.find(f"{{{ns}}}v")
                        cat_labels.append(v.text if v is not None else "")

                result.append({
                    "name":       ser_name,
                    "values":     values,
                    "cat_labels": cat_labels,
                    "ser_el":     ser_el,
                })
    except Exception:
        pass
    return result


def _add_statsig_marker_to_label(shape, chart, ser_idx, cat_idx,
                                  marker, marker_color, slide_num, series_name):
    """
    Append the stat sig marker to the data label text for the given
    series/category point, using green colour.
    Works via python-pptx data_label API.
    """
    try:
        series = list(chart.series)[ser_idx]
        point  = list(series.points)[cat_idx]
        label  = point.data_label
        if label is None or not label.has_text_frame:
            return False
        tf = label.text_frame
        # Append marker as a new run in the last paragraph
        para = tf.paragraphs[-1]
        run  = para.add_run()
        run.text           = f" {marker}"
        run.font.color.rgb = rgb_color(marker_color)
        run.font.bold      = True
        run.font.name      = CHECKLIST["global_font"]
        return True
    except Exception:
        return False


def run_statsig_on_chart(chart, slide, slide_num):
    """
    Main stat sig dispatcher. Returns list of (type, message) change entries.
    Detects mode (products vs others / trend) per chart automatically.
    """
    changes  = []
    mode     = CHECKLIST["statsig_mode"]
    z_crit   = CHECKLIST["statsig_z_critical"]
    marker   = CHECKLIST["statsig_marker"]
    m_color  = CHECKLIST["statsig_marker_color"]

    # Extract all n values visible on the slide
    slide_ns = _extract_all_n_from_slide(slide)

    series_data = _get_series_values(chart)
    if not series_data:
        changes.append(("warn", f"Slide {slide_num}: could not read chart series data"))
        return changes

    # ── Determine comparison mode ─────────────────────────────────
    # Auto-detect: if any series name looks like a quarter → trend mode
    # regardless of user selection, so the right test is always applied
    has_quarter_cats = any(
        _is_quarter_label(lbl)
        for s in series_data
        for lbl in s.get("cat_labels", [])
    )
    has_quarter_series = any(_is_quarter_label(s["name"]) for s in series_data)
    is_trend = has_quarter_cats or has_quarter_series or mode == "Current vs Previous quarter"

    if is_trend:
        changes += _statsig_trend(series_data, chart, slide, slide_num,
                                   z_crit, marker, m_color, slide_ns)
    else:
        changes += _statsig_products_vs_others(series_data, chart, slide, slide_num,
                                                z_crit, marker, m_color, slide_ns)

    return changes


def _find_n_for_label(label_text, slide_ns, fallback_n=None):
    """
    Try to find sample size for a given label by matching against
    the n values extracted from the slide text.
    """
    if not label_text:
        return fallback_n
    lt = label_text.strip()
    # Direct match
    if lt in slide_ns:
        return slide_ns[lt]
    # Partial match — label appears as substring of a footnote line
    for key, n in slide_ns.items():
        if lt.lower() in key.lower() or key.lower() in lt.lower():
            return n
    return fallback_n


def _statsig_products_vs_others(series_data, chart, slide, slide_num,
                                 z_crit, marker, m_color, slide_ns):
    """
    For each series that is NOT 'Others', compare it to the 'Others' series
    at each category. Add marker where stat sig.
    """
    changes = []

    # Find the 'Others' series
    others_idx  = None
    others_data = None
    for i, s in enumerate(series_data):
        if _is_others_series(s["name"]):
            others_idx  = i
            others_data = s
            break

    if others_data is None:
        changes.append(("warn",
            f"Slide {slide_num}: no 'Others' series found — "
            f"skipping product vs others stat sig. "
            f"(Series found: {[s['name'] for s in series_data]})"))
        return changes

    n_others = _find_n_for_label(others_data["name"], slide_ns)

    for ser_idx, s in enumerate(series_data):
        if ser_idx == others_idx:
            continue
        ser_name = s["name"]
        n_prod   = _find_n_for_label(ser_name, slide_ns)

        for cat_idx, val_prod in enumerate(s["values"]):
            if val_prod is None:
                continue
            val_others = (others_data["values"][cat_idx]
                          if cat_idx < len(others_data["values"]) else None)
            if val_others is None:
                continue

            # Get category label for n lookup fallback
            cat_lbl = (s["cat_labels"][cat_idx]
                       if cat_idx < len(s.get("cat_labels", [])) else "")
            n1 = _find_n_for_label(cat_lbl, slide_ns, n_prod)
            n2 = _find_n_for_label(cat_lbl, slide_ns, n_others)

            if not n1 or not n2:
                changes.append(("warn",
                    f"Slide {slide_num} · '{ser_name}' cat {cat_idx+1}: "
                    f"sample size not found — cannot test. "
                    f"Ensure 's = N' appears in chart footnotes."))
                continue

            z = _two_prop_z_test(val_prod, n1, val_others, n2)

            if z >= z_crit:
                ok = _add_statsig_marker_to_label(None, chart, ser_idx, cat_idx,
                                                   marker, m_color, slide_num, ser_name)
                status = "added" if ok else "detected (label not editable)"
                changes.append(("statsig",
                    f"Slide {slide_num} · '{ser_name}' vs Others · "
                    f"Cat {cat_idx+1} ({cat_lbl}): "
                    f"{val_prod:.0f}% vs {val_others:.0f}% · "
                    f"z={z:.2f} ≥ {z_crit} · {marker} {status}"))
            else:
                changes.append(("fix",
                    f"Slide {slide_num} · '{ser_name}' vs Others · "
                    f"Cat {cat_idx+1}: {val_prod:.0f}% vs {val_others:.0f}% · "
                    f"z={z:.2f} — not sig"))

    return changes


def _statsig_trend(series_data, chart, slide, slide_num,
                   z_crit, marker, m_color, slide_ns):
    """
    For each series, compare the most recent time point (last category)
    against the previous one. Add marker where stat sig.

    If categories are on the SERIES axis (e.g. rows = waves), compare
    the first two series instead.
    """
    changes = []

    # Determine whether time is on the category axis or series axis
    # Time on category axis: single series with multiple time-point columns
    # Time on series axis:   multiple series named by quarter

    time_on_series = any(_is_quarter_label(s["name"]) for s in series_data)

    if time_on_series:
        # Series = time points. Compare series[0] (most recent) vs series[1]
        if len(series_data) < 2:
            changes.append(("warn", f"Slide {slide_num}: only one time-series — cannot compare"))
            return changes

        s_curr = series_data[0]
        s_prev = series_data[1]
        n_curr = _find_n_for_label(s_curr["name"], slide_ns)
        n_prev = _find_n_for_label(s_prev["name"], slide_ns)

        for cat_idx, val_curr in enumerate(s_curr["values"]):
            if val_curr is None:
                continue
            val_prev = (s_prev["values"][cat_idx]
                        if cat_idx < len(s_prev["values"]) else None)
            if val_prev is None:
                continue

            cat_lbl = (s_curr["cat_labels"][cat_idx]
                       if cat_idx < len(s_curr.get("cat_labels", [])) else "")
            n1 = _find_n_for_label(cat_lbl, slide_ns, n_curr)
            n2 = _find_n_for_label(cat_lbl, slide_ns, n_prev)

            if not n1 or not n2:
                changes.append(("warn",
                    f"Slide {slide_num} · Trend · Cat {cat_idx+1}: "
                    f"n not found — ensure 's = N' is in footnotes"))
                continue

            z = _two_prop_z_test(val_curr, n1, val_prev, n2)

            if z >= z_crit:
                ok = _add_statsig_marker_to_label(None, chart, 0, cat_idx,
                                                   marker, m_color, slide_num, s_curr["name"])
                status = "added" if ok else "detected (label not editable)"
                changes.append(("statsig",
                    f"Slide {slide_num} · Trend · "
                    f"'{s_curr['name']}' vs '{s_prev['name']}' · "
                    f"Cat {cat_idx+1} ({cat_lbl}): "
                    f"{val_curr:.0f}% vs {val_prev:.0f}% · "
                    f"z={z:.2f} ≥ {z_crit} · {marker} {status}"))

    else:
        # Time on category axis: last 2 categories = current + previous
        for ser_idx, s in enumerate(series_data):
            vals = [v for v in s["values"] if v is not None]
            if len(vals) < 2:
                continue

            # Last category = current, second-to-last = previous
            cat_curr_idx = len(s["values"]) - 1
            cat_prev_idx = cat_curr_idx - 1
            # Walk back to find previous non-None value
            while cat_prev_idx >= 0 and s["values"][cat_prev_idx] is None:
                cat_prev_idx -= 1
            if cat_prev_idx < 0:
                continue

            val_curr = s["values"][cat_curr_idx]
            val_prev = s["values"][cat_prev_idx]

            cats = s.get("cat_labels", [])
            lbl_curr = cats[cat_curr_idx] if cat_curr_idx < len(cats) else f"Cat {cat_curr_idx+1}"
            lbl_prev = cats[cat_prev_idx] if cat_prev_idx < len(cats) else f"Cat {cat_prev_idx+1}"

            n1 = _find_n_for_label(lbl_curr, slide_ns)
            n2 = _find_n_for_label(lbl_prev, slide_ns)

            if not n1 or not n2:
                changes.append(("warn",
                    f"Slide {slide_num} · Trend · '{s['name']}': "
                    f"n not found for '{lbl_curr}' or '{lbl_prev}' — "
                    f"ensure 's = N' appears in chart labels or footnotes"))
                continue

            z = _two_prop_z_test(val_curr, n1, val_prev, n2)

            if z >= z_crit:
                ok = _add_statsig_marker_to_label(None, chart, ser_idx, cat_curr_idx,
                                                   marker, m_color, slide_num, s["name"])
                status = "added" if ok else "detected (label not editable)"
                changes.append(("statsig",
                    f"Slide {slide_num} · Trend · '{s['name']}' · "
                    f"'{lbl_curr}' vs '{lbl_prev}': "
                    f"{val_curr:.0f}% vs {val_prev:.0f}% · "
                    f"z={z:.2f} ≥ {z_crit} · {marker} {status}"))
            else:
                changes.append(("fix",
                    f"Slide {slide_num} · Trend · '{s['name']}' · "
                    f"'{lbl_curr}' vs '{lbl_prev}': "
                    f"{val_curr:.0f}% vs {val_prev:.0f}% · "
                    f"z={z:.2f} — not sig"))

    return changes


# ─────────────────────────────────────────────────────────────────
#  OVERLAP ENGINE
# ─────────────────────────────────────────────────────────────────

def _pts_to_frac(pts):
    return pts / 240.0

def _get_or_create_dlbl(ser_el, idx):
    ns = _NS_C
    for dlbl in ser_el.findall(f"{{{ns}}}dLbl"):
        idx_el = dlbl.find(f"{{{ns}}}idx")
        if idx_el is not None and idx_el.get("val") == str(idx):
            return dlbl
    dlbl    = etree.SubElement(ser_el, f"{{{ns}}}dLbl")
    idx_el  = etree.SubElement(dlbl,   f"{{{ns}}}idx")
    idx_el.set("val", str(idx))
    show    = etree.SubElement(dlbl,   f"{{{ns}}}showVal")
    show.set("val", "1")
    return dlbl

def _apply_nudge(dlbl_el, y_frac):
    ns = _NS_C
    for old in dlbl_el.findall(f"{{{ns}}}layout"):
        dlbl_el.remove(old)
    layout = etree.SubElement(dlbl_el, f"{{{ns}}}layout")
    ml     = etree.SubElement(layout,  f"{{{ns}}}manualLayout")
    for tag, val in [("xMode","factor"),("yMode","factor"),
                     ("x","0.000000"),(  "y", f"{y_frac:.6f}")]:
        el = etree.SubElement(ml, f"{{{ns}}}{tag}")
        el.set("val", val)

def fix_label_overlaps(chart, slide_num, shape_name):
    if not CHECKLIST.get("label_overlap_enabled"):
        return []
    changes   = []
    threshold = CHECKLIST["label_overlap_threshold"]
    nfrac     = _pts_to_frac(CHECKLIST["label_nudge_pts"])
    ns        = _NS_C
    SUPPORTED = {"barChart","bar3DChart","lineChart","line3DChart","areaChart","area3DChart"}
    try:
        plot_el = chart._element.find(f".//{{{ns}}}plotArea")
        if plot_el is None:
            return []
        for ctype in plot_el:
            if etree.QName(ctype.tag).localname not in SUPPORTED:
                continue
            cat_data = {}
            for ser_idx, ser_el in enumerate(ctype.findall(f"{{{ns}}}ser")):
                cache = ser_el.find(f".//{{{ns}}}val/{{{ns}}}numRef/{{{ns}}}numCache")
                if cache is None:
                    cache = ser_el.find(f".//{{{ns}}}numRef/{{{ns}}}numCache")
                if cache is None:
                    continue
                for pt in cache.findall(f"{{{ns}}}pt"):
                    try:
                        idx = int(pt.get("idx", -1))
                        v   = pt.find(f"{{{ns}}}v")
                        if v is None or not v.text:
                            continue
                        val = float(v.text)
                        cat_data.setdefault(idx, []).append((val, ser_idx, ser_el))
                    except (ValueError, TypeError):
                        continue
            for cat_idx, entries in sorted(cat_data.items()):
                if len(entries) < 2:
                    continue
                entries.sort(key=lambda x: x[0], reverse=True)
                for i in range(len(entries) - 1):
                    val_hi, _, ser_hi = entries[i]
                    val_lo, _, ser_lo = entries[i + 1]
                    gap = val_hi - val_lo
                    if gap > threshold:
                        continue
                    try:
                        _apply_nudge(_get_or_create_dlbl(ser_hi, cat_idx), -nfrac)
                        _apply_nudge(_get_or_create_dlbl(ser_lo, cat_idx), +nfrac)
                        changes.append(("overlap",
                            f"Slide {slide_num} '{shape_name}' Cat {cat_idx+1}: "
                            f"{val_hi:.0f}% vs {val_lo:.0f}% (gap {gap:.1f}%) — nudged apart"))
                    except Exception as e:
                        changes.append(("warn",
                            f"Slide {slide_num} '{shape_name}' Cat {cat_idx+1}: "
                            f"overlap could not be fixed ({e})"))
    except Exception as e:
        changes.append(("warn", f"Overlap analysis failed for '{shape_name}': {e}"))
    return changes


# ─────────────────────────────────────────────────────────────────
#  PLACEMENT ENGINE
# ─────────────────────────────────────────────────────────────────

def get_anchors(slide, h):
    a = {"talking": None, "subheader": None, "chart": None}
    for shape in slide.shapes:
        if is_header(shape, h) and not a["talking"]:
            a["talking"] = {"top": shape.top, "left": shape.left}
        elif is_subheader(shape, h) and not a["subheader"]:
            a["subheader"] = {"top": shape.top, "left": shape.left}
    for shape in slide.shapes:
        if shape.has_chart and not a["chart"]:
            a["chart"] = {"top": shape.top, "left": shape.left}
    return a

def apply_placement(slide, h, anchors, slide_num):
    changes = []
    mode    = placement_mode
    do_top  = "top"  in mode.lower() or "+" in mode
    do_left = "left" in mode.lower() or "+" in mode

    def snap(shape, anchor, label):
        moved = []
        if do_top and anchor.get("top") is not None:
            diff = abs((shape.top or 0) - anchor["top"])
            if diff > 0:
                shape.top = anchor["top"]
                moved.append(f"top by {diff//914:.2f}\"")
        if do_left and anchor.get("left") is not None:
            diff = abs((shape.left or 0) - anchor["left"])
            if diff > 0:
                shape.left = anchor["left"]
                moved.append(f"left by {diff//914:.2f}\"")
        if moved:
            changes.append(("fix",
                f"Slide {slide_num} · {label} snapped ({', '.join(moved)})"))

    for shape in slide.shapes:
        if align_talking   and anchors["talking"]   and is_header(shape, h):
            snap(shape, anchors["talking"],   "talking header")
        elif align_subheader and anchors["subheader"] and is_subheader(shape, h):
            snap(shape, anchors["subheader"], "sub header")
        elif align_chart     and anchors["chart"]     and shape.has_chart:
            snap(shape, anchors["chart"],     "chart")
    return changes


# ─────────────────────────────────────────────────────────────────
#  MAIN SLIDE PROCESSOR
# ─────────────────────────────────────────────────────────────────

def process_slide(slide, h, slide_num, anchors):
    changes = []

    # ── Text shapes ───────────────────────────────────────────────
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if is_footnote(shape, h):
            for para in tf.paragraphs:
                for run in para.runs:
                    if CHECKLIST["footnote_remove_highlight"]:
                        remove_highlight(run)
                    if is_low_sample(run.text):
                        set_run(run, size=CHECKLIST["footnote_size"],
                                color=CHECKLIST["footnote_low_sample_color"],
                                font=CHECKLIST["global_font"])
                        changes.append(("fix", f"Slide {slide_num} · Footnote low-sample red"))
                    else:
                        set_run(run, size=CHECKLIST["footnote_size"],
                                color=CHECKLIST["footnote_color"],
                                font=CHECKLIST["global_font"])
                        changes.append(("fix", f"Slide {slide_num} · Footnote grey {CHECKLIST['footnote_size']}pt"))
        elif is_header(shape, h):
            for para in tf.paragraphs:
                for run in para.runs:
                    set_run(run, size=CHECKLIST["talking_header_size"],
                            color=CHECKLIST["talking_header_color"],
                            font=CHECKLIST["global_font"])
            changes.append(("fix", f"Slide {slide_num} · Header {CHECKLIST['talking_header_size']}pt navy"))
        elif is_subheader(shape, h):
            for i, para in enumerate(tf.paragraphs):
                sz = CHECKLIST["subheader_size"] if i == 0 else CHECKLIST["below_subheader_size"]
                for run in para.runs:
                    set_run(run, size=sz, font=CHECKLIST["global_font"])
            changes.append(("fix", f"Slide {slide_num} · Sub header formatted"))
        else:
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = CHECKLIST["global_font"]

    # ── Charts ────────────────────────────────────────────────────
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart      = shape.chart
        shape_name = shape.name or "Chart"

        # Chart title
        if chart.has_title and chart.chart_title.has_text_frame:
            for para in chart.chart_title.text_frame.paragraphs:
                for run in para.runs:
                    set_run(run, size=CHECKLIST["chart_title_size"],
                            bold=CHECKLIST["chart_title_bold"],
                            color=CHECKLIST["chart_title_color"],
                            font=CHECKLIST["global_font"])
            changes.append(("fix", f"Slide {slide_num} · Chart title formatted"))

        # Data labels
        for series in chart.series:
            try:
                fg = series.format.fill.fore_color.rgb
                r, g, b  = fg[0], fg[1], fg[2]
                lbl_col  = (CHECKLIST["data_label_dark_color"]
                             if is_dark_bar(r, g, b)
                             else CHECKLIST["data_label_light_color"])
            except Exception:
                lbl_col  = CHECKLIST["data_label_light_color"]
            try:
                for point in series.points:
                    try:
                        label = point.data_label
                        if label is None:
                            continue
                        label_text = label.text_frame.text if label.has_text_frame else ""
                        pct = pct_value(label_text)
                        if pct is not None and pct < CHECKLIST["data_label_min_pct"]:
                            if label.has_text_frame:
                                for para in label.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                            changes.append(("fix", f"Slide {slide_num} · Removed {pct:.0f}% label"))
                            continue
                        if label.has_text_frame:
                            for para in label.text_frame.paragraphs:
                                for run in para.runs:
                                    set_run(run, font=CHECKLIST["global_font"], color=lbl_col)
                                    if CHECKLIST["data_label_remove_highlight"]:
                                        remove_highlight(run)
                    except Exception:
                        continue
            except Exception:
                pass

        # Axis labels
        try:
            for axis in [chart.category_axis, chart.value_axis]:
                try:
                    axis.tick_labels.font.name = CHECKLIST["global_font"]
                except Exception:
                    pass
        except Exception:
            pass

        # Overlap fix
        changes.extend(fix_label_overlaps(chart, slide_num, shape_name))

        # Stat sig
        if CHECKLIST["statsig_enabled"]:
            allowed = CHECKLIST["statsig_slides"]
            if allowed is None or slide_num in allowed:
                changes.extend(run_statsig_on_chart(chart, slide, slide_num))

    # ── Placement ─────────────────────────────────────────────────
    if placement_enabled and anchors and slide_num > 1:
        changes.extend(apply_placement(slide, h, anchors, slide_num))

    return changes


# ─────────────────────────────────────────────────────────────────
#  RUN FORMATTER
# ─────────────────────────────────────────────────────────────────

def run_formatter(file_bytes):
    prs     = Presentation(io.BytesIO(file_bytes))
    h       = prs.slide_height
    all_chg = []
    anchors = {}

    if placement_enabled and prs.slides:
        anchors = get_anchors(prs.slides[0], h)

    for i, slide in enumerate(prs.slides, start=1):
        all_chg.extend(process_slide(slide, h, i, anchors))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue(), all_chg


# ─────────────────────────────────────────────────────────────────
#  MAIN UI
# ─────────────────────────────────────────────────────────────────

st.title("🎯 Deck Formatting Agent")
st.markdown(
    "Fill in your client rules in the **sidebar**, upload your deck below, "
    "click **Run**. Get a corrected deck back instantly."
)
st.markdown("---")

col_up, col_btn = st.columns([3, 1])
with col_up:
    uploaded = st.file_uploader(
        "Upload your .pptx deck",
        type=["pptx"],
        help="Processed entirely in memory — nothing is stored.",
    )
with col_btn:
    st.markdown("<br>", unsafe_allow_html=True)
    run_btn = st.button("✨ Run Formatter", type="primary",
                        use_container_width=True, disabled=uploaded is None)

if uploaded is None:
    st.info("👈  Set your rules in the sidebar, then upload your deck above.")
    st.stop()

st.success(f"📎 **{uploaded.name}** ready")

if run_btn:
    with st.spinner("Analysing and correcting your deck..."):
        try:
            corrected_bytes, changes = run_formatter(uploaded.read())

            fixes    = [c for c in changes if c[0] == "fix"]
            statsigs = [c for c in changes if c[0] == "statsig"]
            overlaps = [c for c in changes if c[0] == "overlap"]
            warnings = [c for c in changes if c[0] == "warn"]

            st.markdown("---")

            # Stats row
            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("Formatting fixes",  len(fixes))
            c2.metric("Stat sig markers",  len(statsigs))
            c3.metric("Overlaps fixed",    len(overlaps))
            c4.metric("Warnings",          len(warnings))
            c5.metric("Status", "✅ Clean" if not warnings else "⚠️ Review")

            st.markdown("---")

            # Download
            st.download_button(
                label="⬇️  Download corrected deck",
                data=corrected_bytes,
                file_name=("corrected_" + (uploaded.name[:-5] if uploaded.name.lower().endswith(".pptx") else uploaded.name) + ".pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )

            st.markdown("---")

            # Change log
            with st.expander(
                f"📋  Full change log  ({len(changes)} items — "
                f"{len(statsigs)} stat sig · {len(overlaps)} overlap · "
                f"{len(fixes)} formatting · {len(warnings)} warnings)",
                expanded=False,
            ):
                if statsigs:
                    st.markdown("**🟢 Stat sig markers added**")
                    for _, msg in statsigs:
                        st.markdown(f'<div class="statsig-item">● {msg}</div>',
                                    unsafe_allow_html=True)

                if overlaps:
                    st.markdown("**↕ Overlapping labels nudged**")
                    for _, msg in overlaps:
                        st.markdown(f'<div class="overlap-item">↕ {msg}</div>',
                                    unsafe_allow_html=True)

                if fixes:
                    st.markdown("**✓ Formatting corrections**")
                    for _, msg in fixes:
                        st.markdown(f'<div class="fix-item">✓ {msg}</div>',
                                    unsafe_allow_html=True)

                if warnings:
                    st.markdown("**⚠ Needs manual review**")
                    for _, msg in warnings:
                        st.markdown(f'<div class="warn-item">⚠ {msg}</div>',
                                    unsafe_allow_html=True)

                if not changes:
                    st.info("No changes needed — deck already meets all rules.")

        except Exception as e:
            st.error(f"Error processing deck: {e}")
            st.info("Make sure your file is a valid .pptx (not .ppt or password-protected).")

elif uploaded:
    st.info("Rules are set. Click **✨ Run Formatter** when ready.")
